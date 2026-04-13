"""
Generate slide deck (PPTX) and talk track (PDF) for Code Apps Tech Series.
Newbie-friendly explanations of advanced findings. Companion to a colleague's intro talk.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from fpdf import FPDF
import os

OUT_DIR = os.path.dirname(os.path.abspath(__file__))

# ── Colors ──────────────────────────────────────────────────────────
PURPLE = RGBColor(0x74, 0x27, 0x74)       # Power Platform purple
DARK_BG = RGBColor(0x1E, 0x1E, 0x2E)      # Dark slide background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_BLUE = RGBColor(0x60, 0xA5, 0xFA)
ACCENT_AMBER = RGBColor(0xFB, 0xBF, 0x24)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
ACCENT_VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
MID_GRAY = RGBColor(0x94, 0xA3, 0xB8)

# ── Helpers ─────────────────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, font_name="Segoe UI", spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing

        # Support (text, accent_color) tuples or plain strings
        if isinstance(item, tuple):
            txt, accent = item
            # Bold prefix before " -- "
            if " -- " in txt:
                prefix, rest = txt.split(" -- ", 1)
                run1 = p.add_run()
                run1.text = "\u2022  " + prefix + " "
                run1.font.size = Pt(font_size)
                run1.font.color.rgb = accent
                run1.font.bold = True
                run1.font.name = font_name
                run2 = p.add_run()
                run2.text = rest
                run2.font.size = Pt(font_size)
                run2.font.color.rgb = color
                run2.font.name = font_name
            else:
                run = p.add_run()
                run.text = "\u2022  " + txt
                run.font.size = Pt(font_size)
                run.font.color.rgb = accent
                run.font.name = font_name
        else:
            run = p.add_run()
            run.text = "\u2022  " + item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = font_name
    return tf

def add_accent_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_code_box(slide, left, top, width, height, text, font_size=11):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = ACCENT_GREEN
    p.font.name = "Consolas"
    return tf

# ── PPTX Generation ────────────────────────────────────────────────
def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 1 — TITLE
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, 13.333, 0.08, PURPLE)

    add_text_box(slide, 0.8, 1.5, 11.7, 1.2,
                 "Code Apps: Under the Hood",
                 font_size=44, color=WHITE, bold=True)
    add_text_box(slide, 0.8, 2.8, 11.7, 0.8,
                 "What we found when we looked inside \u2014 and what it means for you",
                 font_size=22, color=LIGHT_GRAY)
    add_accent_bar(slide, 0.8, 3.8, 3.0, 0.04, ACCENT_VIOLET)
    add_text_box(slide, 0.8, 4.2, 6, 0.6,
                 "Power Platform Tech Series  \u00b7  2026",
                 font_size=16, color=MID_GRAY)

    # Tag line at bottom
    add_text_box(slide, 0.8, 6.2, 11.7, 0.5,
                 "Companion to \"Intro to Code Apps\" \u2014 we pick up where that leaves off",
                 font_size=14, color=MID_GRAY, alignment=PP_ALIGN.LEFT)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 2 — YOUR APP LIVES INSIDE A BIGGER SYSTEM
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, 13.333, 0.06, PURPLE)

    add_text_box(slide, 0.8, 0.4, 11.7, 0.8,
                 "Your App Lives Inside a Bigger System",
                 font_size=36, color=WHITE, bold=True)
    add_text_box(slide, 0.8, 1.1, 11.7, 0.5,
                 "Think of it like renting space in an office building \u2014 you get security, elevators, and fire alarms for free",
                 font_size=16, color=MID_GRAY)

    # Left column — what the platform does behind the scenes
    add_text_box(slide, 0.8, 1.8, 5.5, 0.5,
                 "What Happens Behind the Scenes", font_size=20, color=ACCENT_VIOLET, bold=True)
    add_bullet_list(slide, 0.8, 2.3, 5.5, 4.5, [
        ("Your app is treated like a Canvas App -- Even though you wrote React code, the platform classifies it the same as a drag-and-drop app", ACCENT_RED),
        ("It\u2019s nested inside multiple layers -- Like a series of picture frames: platform shell \u2192 player \u2192 runtime \u2192 your code", ACCENT_AMBER),
        ("77 scripts load before yours -- For login, security, Teams integration, Copilot AI, user surveys, and more", ACCENT_BLUE),
        ("230KB of translations load every time -- Text for every language, device type, and region \u2014 most of it will never be used", ACCENT_GREEN),
        ("55 on/off switches -- Microsoft can enable or disable features for your organization remotely", MID_GRAY),
    ], font_size=14, spacing=Pt(8))

    # Right column — what you get without doing anything
    add_text_box(slide, 7.2, 1.8, 5.5, 0.5,
                 "What You Get Without Doing Anything", font_size=20, color=ACCENT_GREEN, bold=True)
    add_bullet_list(slide, 7.2, 2.3, 5.5, 4.5, [
        ("Login is handled for you -- Users are authenticated before your app even loads \u2014 you never touch passwords or tokens", ACCENT_GREEN),
        ("Copilot AI is pre-wired -- A built-in AI assistant panel is ready to go \u2014 185 references to it in the page code", ACCENT_VIOLET),
        ("Works in government clouds automatically -- Same app code runs in US Gov, Germany, China \u2014 no changes needed", ACCENT_BLUE),
        ("Accessibility is built in -- Screen reader support, keyboard navigation, skip-to-content links, progress indicators", ACCENT_AMBER),
        ("Security headers are managed -- Content security policies and script protections are applied for you", MID_GRAY),
    ], font_size=14, spacing=Pt(8))

    # Bottom callout
    add_accent_bar(slide, 0.8, 6.6, 11.7, 0.04, PURPLE)
    add_text_box(slide, 0.8, 6.7, 11.7, 0.5,
                 "How we know: we saved the deployed app\u2019s HTML and read every line \u2014 I\u2019ll show you in the demo",
                 font_size=13, color=MID_GRAY)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 3 — SURPRISES WE FOUND
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, 13.333, 0.06, ACCENT_RED)

    add_text_box(slide, 0.8, 0.4, 11.7, 0.8,
                 "Surprises We Found",
                 font_size=36, color=WHITE, bold=True)
    add_text_box(slide, 0.8, 1.1, 11.7, 0.5,
                 "Two real bugs that looked like everything was working fine",
                 font_size=16, color=MID_GRAY)

    # Left column — Story 1: The invisible data bug
    add_accent_bar(slide, 0.8, 1.85, 0.06, 4.3, ACCENT_AMBER)
    add_text_box(slide, 1.1, 1.8, 5.3, 0.5,
                 "\u26a0  The Invisible Data Bug", font_size=20, color=ACCENT_AMBER, bold=True)
    add_text_box(slide, 1.1, 2.3, 5.3, 0.4,
                 "Saving worked perfectly \u2014 but opening the same record showed the link was \"gone\"",
                 font_size=13, color=LIGHT_GRAY)

    add_bullet_list(slide, 1.1, 2.8, 5.3, 3.5, [
        ("The data was actually fine -- The contact was linked to the right account in the database. The app just couldn\u2019t read it back.", ACCENT_AMBER),
        ("Why? Different field names for saving vs. reading -- The system uses one name when you write data and a completely different name when you read it", LIGHT_GRAY),
        ("It looked like it worked -- You could save, close, and see success. Only when you reopened the record did you notice the link showed \"None\"", LIGHT_GRAY),
        ("The fix: check both field names -- A simple fallback \u2014 look for the data under both names and use whichever one has it", ACCENT_GREEN),
    ], font_size=13, spacing=Pt(6))

    # Right column — Story 2: The chatbot that only worked locally
    add_accent_bar(slide, 7.0, 1.85, 0.06, 4.3, ACCENT_RED)
    add_text_box(slide, 7.3, 1.8, 5.3, 0.5,
                 "\u274c  The Chatbot That Only Worked Locally", font_size=20, color=ACCENT_RED, bold=True)
    add_text_box(slide, 7.3, 2.3, 5.3, 0.4,
                 "Worked great on my machine \u2014 completely broke when deployed",
                 font_size=13, color=LIGHT_GRAY)

    add_bullet_list(slide, 7.3, 2.8, 5.3, 3.5, [
        ("We copied a proven pattern -- An AI chatbot setup that worked perfectly in a standalone website", LIGHT_GRAY),
        ("Code Apps handle login differently -- The platform uses its own login system, so our chatbot couldn\u2019t get the credentials it needed", ACCENT_RED),
        ("Local testing hid the problem -- The dev environment doesn\u2019t show the same login restrictions as the real deployment", LIGHT_GRAY),
        ("The fix: use the platform\u2019s built-in chat -- Since both the app and chatbot live in Power Platform, the login session is already shared", ACCENT_GREEN),
    ], font_size=13, spacing=Pt(6))

    # Bottom callout
    add_accent_bar(slide, 0.8, 6.5, 11.7, 0.04, ACCENT_RED)
    add_text_box(slide, 0.8, 6.6, 11.7, 0.6,
                 "Lesson: always test the deployed version \u2014 what works locally can fail in production",
                 font_size=15, color=ACCENT_AMBER, bold=True, alignment=PP_ALIGN.CENTER)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 4 — AN AI AGENT BUILT THIS ENTIRE APP
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, 13.333, 0.06, ACCENT_GREEN)

    add_text_box(slide, 0.8, 0.4, 11.7, 0.8,
                 "An AI Agent Built This Entire App",
                 font_size=36, color=WHITE, bold=True)
    add_text_box(slide, 0.8, 1.1, 11.7, 0.5,
                 "18 rounds of work \u2014 from blank project to full app \u2014 every decision is documented",
                 font_size=16, color=MID_GRAY)

    # Left — the build story
    add_text_box(slide, 0.8, 1.8, 5.5, 0.5,
                 "How We Built It", font_size=20, color=ACCENT_GREEN, bold=True)
    add_bullet_list(slide, 0.8, 2.3, 5.5, 2.5, [
        ("An AI assistant wrote every line -- Scaffolding, data screens, UI, AI features, drag-and-drop board \u2014 all of it", ACCENT_GREEN),
        ("Each round: research first, then build -- \"How does this work in Code Apps?\" \u2192 discover the gotchas \u2192 build the feature", ACCENT_BLUE),
        ("Every decision is written down -- A full journal of what was tried, what broke, and how it was fixed", ACCENT_AMBER),
        ("The AI hit real bugs -- The same tricky platform behaviors that would trip up a human developer", ACCENT_RED),
    ], font_size=14, spacing=Pt(8))

    # Right — the dashboard story
    add_text_box(slide, 7.2, 1.8, 5.5, 0.5,
                 "The Dashboard: Simpler Than It Looks", font_size=20, color=ACCENT_VIOLET, bold=True)
    add_bullet_list(slide, 7.2, 2.3, 5.5, 2.5, [
        ("705 lines, no charting library -- The charts are drawn with basic web graphics \u2014 circles and rectangles", ACCENT_VIOLET),
        ("One prompt rebuilt it from scratch -- \"Replace the chart library with something simpler\" \u2192 done", ACCENT_GREEN),
        ("Every chart is clickable -- Click any number, bar, or ring segment to see the actual records behind it", ACCENT_BLUE),
        ("Kanban board: drag cards between columns -- 4 columns with real-time saves to the database", ACCENT_AMBER),
    ], font_size=14, spacing=Pt(8))

    # Bottom — the key stats
    add_accent_bar(slide, 0.8, 5.0, 11.7, 0.04, ACCENT_GREEN)

    # Stats row
    stats = [
        ("6", "Database\nTables"),
        ("18", "Build\nRounds"),
        ("705", "Dashboard\nLines of Code"),
        ("0", "Chart\nLibraries"),
        ("1", "Deploy\nCommand"),
    ]
    stat_width = 2.0
    stat_start = 1.2
    for i, (num, label) in enumerate(stats):
        x = stat_start + i * 2.3
        add_text_box(slide, x, 5.3, stat_width, 0.6,
                     num, font_size=36, color=ACCENT_GREEN, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name="Consolas")
        add_text_box(slide, x, 5.9, stat_width, 0.6,
                     label, font_size=12, color=MID_GRAY,
                     alignment=PP_ALIGN.CENTER)

    add_text_box(slide, 0.8, 6.7, 11.7, 0.5,
                 "npm run build && pac code push  \u2014  that's the entire deploy",
                 font_size=14, color=ACCENT_GREEN, font_name="Consolas",
                 alignment=PP_ALIGN.LEFT)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 5 — COOL THINGS YOU CAN DO
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, 13.333, 0.06, ACCENT_BLUE)

    add_text_box(slide, 0.8, 0.4, 11.7, 0.8,
                 "Cool Things You Can Do",
                 font_size=36, color=WHITE, bold=True)
    add_text_box(slide, 0.8, 1.1, 11.7, 0.5,
                 "Standard tools that just work inside Power Platform \u2014 no special tricks needed",
                 font_size=16, color=MID_GRAY)

    # Four pattern cards as two rows of two
    patterns = [
        ("Instant Search, Zero Extra Calls", ACCENT_BLUE,
         "Press Ctrl+K to search all 6 data tables\ninstantly. It reads data already loaded\nin memory \u2014 no extra database trips.",
         "Searches accounts, contacts, action items, ideas, meetings, projects"),
        ("Quick-Create Buttons (27 lines of code)", ACCENT_AMBER,
         "Colored buttons at the top of every page.\nTap one and the right form opens,\npre-filled with the right type. That\u2019s it.",
         "Tap \"work\" \u2192 action item form opens with type = Work"),
        ("Drag-and-Drop That Saves Automatically", ACCENT_GREEN,
         "Drag a card on the Kanban board to a\nnew column and the database updates\ninstantly. No save button needed.",
         "Drag to parking lot \u2192 item is flagged as pinned in Dataverse"),
        ("Dark Mode That Works Inside the Platform", ACCENT_VIOLET,
         "One click toggles dark mode. This needed\na specific approach because the platform\ncontrols the outer page styling.",
         "Class-based toggle \u2014 media queries don\u2019t work inside the iframe"),
    ]

    for i, (title, accent, desc, code) in enumerate(patterns):
        col = i % 2
        row = i // 2
        x = 0.8 + col * 6.2
        y = 1.8 + row * 2.6

        add_accent_bar(slide, x, y, 0.06, 2.2, accent)
        add_text_box(slide, x + 0.25, y + 0.05, 5.5, 0.4,
                     title, font_size=18, color=accent, bold=True)
        add_text_box(slide, x + 0.25, y + 0.5, 5.5, 1.1,
                     desc, font_size=13, color=LIGHT_GRAY)
        add_code_box(slide, x + 0.25, y + 1.6, 5.5, 0.4,
                     code, font_size=10)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # SLIDE 6 — LET ME SHOW YOU (demo transition)
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_accent_bar(slide, 0, 0, 13.333, 0.06, ACCENT_AMBER)

    add_text_box(slide, 0.8, 1.0, 11.7, 1.0,
                 "Let Me Show You",
                 font_size=48, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_accent_bar(slide, 5.5, 2.3, 2.3, 0.04, ACCENT_AMBER)
    add_text_box(slide, 0.8, 2.7, 11.7, 0.6,
                 "Live Demo",
                 font_size=28, color=ACCENT_AMBER,
                 alignment=PP_ALIGN.CENTER)

    # Demo flow bullets — centered
    demo_items = [
        "\u2776  Dashboard \u2014 click any number to see the records behind it",
        "\u2777  Kanban Board \u2014 drag cards between columns, watch the database update",
        "\u2778  Ctrl+K Search \u2014 find anything across all data tables instantly",
        "\u2779  AI Magic \u2014 paste meeting notes, watch action items get created automatically",
        "\u277a  Dark Mode \u2014 one click, works inside the platform",
        "\u277b  Under the Hood \u2014 open dev tools to see the 77 scripts and the runtime",
        "\u277c  Deploy \u2014 two commands and it\u2019s live",
    ]
    tf = add_bullet_list(slide, 2.5, 3.5, 8.3, 3.5, demo_items,
                         font_size=16, color=LIGHT_GRAY, spacing=Pt(6))

    pptx_path = os.path.join(OUT_DIR, "code-apps-under-the-hood.pptx")
    prs.save(pptx_path)
    print(f"PPTX saved: {pptx_path}")
    return pptx_path


# ── PDF Talk Track ──────────────────────────────────────────────────
def build_pdf():
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=20)
    pdf.set_margins(25, 20, 25)

    # ── Fonts ───────────────────────────────────────────────────────
    pdf.add_font("Segoe", "", "C:/Windows/Fonts/segoeui.ttf", uni=True)
    pdf.add_font("Segoe", "B", "C:/Windows/Fonts/segoeuib.ttf", uni=True)
    pdf.add_font("Segoe", "I", "C:/Windows/Fonts/segoeuii.ttf", uni=True)
    pdf.add_font("Consolas", "", "C:/Windows/Fonts/consola.ttf", uni=True)

    def heading(text, size=18):
        pdf.set_font("Segoe", "B", size)
        pdf.set_text_color(116, 39, 116)  # Purple
        pdf.cell(0, 10, text, new_x="LMARGIN", new_y="NEXT")
        pdf.ln(2)

    def subheading(text, size=13):
        pdf.set_font("Segoe", "B", size)
        pdf.set_text_color(40, 40, 60)
        pdf.cell(0, 8, text, new_x="LMARGIN", new_y="NEXT")
        pdf.ln(1)

    def body(text, size=10.5):
        pdf.set_font("Segoe", "", size)
        pdf.set_text_color(50, 50, 50)
        pdf.multi_cell(0, 5.5, text)
        pdf.ln(2)

    def stage_direction(text, size=10):
        pdf.set_font("Segoe", "I", size)
        pdf.set_text_color(120, 120, 140)
        pdf.multi_cell(0, 5, text)
        pdf.ln(2)

    def code(text, size=9):
        pdf.set_font("Consolas", "", size)
        pdf.set_text_color(16, 185, 129)  # Green
        pdf.multi_cell(0, 5, text)
        pdf.set_font("Segoe", "", 10.5)
        pdf.set_text_color(50, 50, 50)
        pdf.ln(2)

    def bullet(text, size=10.5):
        pdf.set_font("Segoe", "", size)
        pdf.set_text_color(50, 50, 50)
        x = pdf.get_x()
        pdf.cell(6, 5.5, "\u2022")
        pdf.multi_cell(0, 5.5, text)
        pdf.ln(1)

    def divider():
        pdf.ln(3)
        pdf.set_draw_color(116, 39, 116)
        pdf.set_line_width(0.3)
        y = pdf.get_y()
        pdf.line(25, y, 185, y)
        pdf.ln(5)

    # ════════════════════════════════════════════════════════════════
    # PAGE 1 — TITLE + OVERVIEW
    # ════════════════════════════════════════════════════════════════
    pdf.add_page()
    pdf.set_font("Segoe", "B", 24)
    pdf.set_text_color(116, 39, 116)
    pdf.cell(0, 14, "Code Apps: Under the Hood", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Segoe", "", 13)
    pdf.set_text_color(100, 100, 100)
    pdf.cell(0, 8, "Talk Track & Speaker Notes", new_x="LMARGIN", new_y="NEXT")
    pdf.cell(0, 8, "Power Platform Tech Series  |  2026", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(6)

    body("This talk track accompanies 6 slides. It follows a colleague's \"Intro to Code Apps\" "
         "presentation, so don't repeat the basics (what Code Apps are, how to scaffold, the CLI). "
         "Instead, focus on what we discovered when we actually built and deployed a real app, "
         "then transition into a live demo.")
    pdf.ln(2)

    subheading("Session Flow")
    bullet("Slides: ~10 minutes (6 slides)")
    bullet("Live demo: ~15-20 minutes")
    bullet("Q&A: remaining time")
    pdf.ln(2)

    subheading("Key Messages (Plain English)")
    bullet("Your Code App doesn't run alone -- it lives inside a massive system that handles login, security, AI, accessibility, and government cloud compliance for you. You just write your app.")
    bullet("This entire demo app was built by an AI coding assistant across 18 rounds of work. Every decision and mistake is documented.")
    bullet("The analytics dashboard has no chart library -- it's 705 lines of basic web graphics. One prompt to the AI rebuilt the whole thing.")
    bullet("Standard web development tools work perfectly inside Power Platform. The platform doesn't get in your way.")
    bullet("Testing locally is not enough. We found two bugs that only appeared after deploying to the real platform.")

    divider()

    # ════════════════════════════════════════════════════════════════
    # SLIDE 1 TALK TRACK
    # ════════════════════════════════════════════════════════════════
    heading("Slide 1: Title")
    stage_direction("[Slide on screen. Brief pause. Let people read it.]")
    body("\"[Colleague] just walked you through what Code Apps are and how to get started. "
         "I'm going to show you what we found when we actually built one. We'll look at what's "
         "really happening behind the scenes when your app runs, a couple of surprises we hit, "
         "how an AI assistant built the whole thing, and then I'll show you the app live.\"")

    divider()

    # ════════════════════════════════════════════════════════════════
    # SLIDE 2 TALK TRACK
    # ════════════════════════════════════════════════════════════════
    heading("Slide 2: Your App Lives Inside a Bigger System")
    stage_direction("[Advance slide. Give people 5 seconds to scan the two columns.]")

    body("\"Here's something that surprised me. After deploying the app, I opened browser dev tools "
         "and saved the entire page HTML. The file was enormous -- 81,000 tokens worth of code. "
         "And my app? Not a single reference to it. No React, no Tailwind, none of my components. "
         "My app was buried deep inside this massive system.\"")
    pdf.ln(1)

    body("\"Think of it like renting space in an office building. You decorate your office however you want, "
         "but the building provides security guards, elevators, fire alarms, and wheelchair ramps. "
         "That's what Power Platform does for your Code App.\"")
    pdf.ln(1)

    subheading("Left Column -- What Happens Behind the Scenes")
    bullet("Your app is classified as a \"ClassicCanvasApp\" -- the same category as a drag-and-drop canvas app. The platform doesn't distinguish between them.")
    bullet("Your app is nested inside multiple layers -- like a series of picture frames. The platform's shell loads a player, the player loads a runtime, and the runtime finally loads your code. Each layer adds capabilities.")
    bullet("77 separate scripts are registered to load -- these handle login (MSAL), legacy compatibility (ADAL, WinJS), security (DOMPurify), Teams integration, the Copilot AI panel, user feedback surveys, and more. All of this loads before your first line of code runs.")
    bullet("230KB of translation text is embedded in every page load -- strings for every language, for barcode scanners, NFC readers, mobile device features, and region-specific compliance (China, Germany, etc.). Most of it will never be used in a browser, but it's always there.")
    bullet("55 feature switches let Microsoft toggle capabilities per-tenant -- things like the Copilot sidebar, newer login protocols, feedback surveys, and security reporting. Your app runs inside a system that's constantly being reconfigured around you.")
    pdf.ln(1)

    subheading("Right Column -- What You Get Without Doing Anything")
    bullet("Login is completely handled for you. Users are authenticated by Microsoft Entra ID before your page even loads. There are no login tokens in the HTML. You never write login code.")
    bullet("The Copilot AI assistant has its own dedicated panel (320 pixels wide), its own scripts, teaching bubbles, and the ability to search records, summarize data, and translate natural language into filters. We counted 185 references to it in the page HTML alone.")
    bullet("Government and sovereign cloud routing is built in. The same app code automatically works in US Government (GCC), Germany, China, and commercial environments. The platform handles the routing -- you change nothing.")
    bullet("Accessibility features are pre-built: screen reader roles, keyboard navigation regions, skip-to-content links, and a multi-stage progress bar with descriptive labels. The hard parts of accessibility are done for you.")
    pdf.ln(1)

    body("\"The takeaway: when you deploy a Code App, you're not just hosting a website. "
         "You're placing your app inside a managed enterprise system that handles login, "
         "security, AI, accessibility, and government compliance. You just write your app.\"")

    stage_direction("[Pause. Let that sink in. If someone asks about performance: acknowledge "
                    "that yes, loading 77 scripts and 230KB of text is overhead -- but it's the same "
                    "overhead every Canvas App already pays. You're not adding to it.]")

    divider()

    # ════════════════════════════════════════════════════════════════
    # SLIDE 3 TALK TRACK
    # ════════════════════════════════════════════════════════════════
    heading("Slide 3: Surprises We Found")
    stage_direction("[Advance slide. Two stories side by side.]")

    body("\"Now that you know what's running behind the scenes, let me tell you about two bugs "
         "we found. Both of them looked like everything was working perfectly -- until they weren't.\"")
    pdf.ln(1)

    subheading("Left Column -- The Invisible Data Bug")
    body("\"We built a contact form with an account dropdown. You pick a company, link it to "
         "a contact, hit save -- great, it works. But when you open that same contact again, "
         "the company field says 'None.' The link is gone. Except... it's not gone.\"")
    pdf.ln(1)
    bullet("The data was actually correct in the database. The link was saved properly. The problem was that the system uses one field name when you SAVE data (parentcustomerid) and a completely different field name when you READ it back (_parentcustomerid_value). The auto-generated code only knew about the save name.")
    bullet("So saving worked perfectly every time. But reading failed silently -- the app looked for the data under the wrong name, found nothing, and showed 'None.' No error message, no warning. It just looked broken.")
    bullet("The fix was simple once we understood it: check both field names. Try the official name first, and if it's empty, check the alternate name. Same thing for the company's display name -- we had to build our own lookup because the auto-generated name field doesn't work reliably either.")
    pdf.ln(1)
    body("\"The lesson: the auto-generated code gives you a contract for saving data, but reading "
         "it back can use completely different names. You have to verify that reads actually work.\"")
    pdf.ln(1)

    subheading("Right Column -- The Chatbot That Only Worked Locally")
    body("\"Second story. We wanted to add an AI chatbot to the app. I had a pattern that "
         "worked great in another project -- a regular Azure website with a Copilot Studio bot. "
         "We used Microsoft's standard login flow to connect the user to the bot. Proven, tested, "
         "production-ready.\"")
    pdf.ln(1)
    bullet("We copied that pattern into the Code App. Locally, it looked great. But here's the thing: Code Apps don't handle login the way regular websites do. The Power Platform uses its own custom login system (called paauth and dynamicauth). There's no way to extract a standard login token from it. Our chatbot connection literally cannot work.")
    bullet("Local testing hid the problem. When you develop locally with 'pac code run,' you're technically inside a Power Platform wrapper, but the login context is different enough that you don't see the failure. You only discover this after deploying to the real environment.")
    bullet("The fix was a complete change of approach. Instead of our custom chatbot setup with three extra code libraries, we just opened the Copilot Studio hosted chat in a simple embedded frame. Since the app and the chatbot both live inside Power Platform, the user's login session is already shared. Zero extra code. Zero extra libraries. It just works.")
    pdf.ln(1)
    body("\"The lesson: patterns that work in regular web apps don't always transfer to Code Apps. "
         "The login model is fundamentally different. Always test the deployed version, not just local.\"")

    stage_direction("[This slide often generates questions. Common ones:\n"
                    "  'Can you use standard login libraries at all?' -- No, not in the traditional sense. The platform owns the login context.\n"
                    "  'Does the embedded chatbot work in all browsers?' -- Yes, same session cookies.\n"
                    "Keep answers brief -- you'll show the working chatbot in the demo.]")

    divider()

    # ════════════════════════════════════════════════════════════════
    # SLIDE 4 TALK TRACK
    # ════════════════════════════════════════════════════════════════
    heading("Slide 4: An AI Agent Built This Entire App")
    stage_direction("[Advance slide. The stats row at the bottom will draw eyes -- let people see the numbers.]")

    body("\"This is the part everyone asks about. This entire app -- six database tables, "
         "full create/read/update/delete on all of them, an analytics dashboard, a Kanban board "
         "with drag-and-drop, a search bar, and AI-powered features -- was built entirely by "
         "an AI coding assistant called Claude Code. Eighteen rounds of work.\"")
    pdf.ln(1)

    subheading("Left Column -- How We Built It")
    bullet("Every round started with a research question. 'How do data relationships work in Code Apps?' The AI would explore the platform, find the tricky parts, and then build the feature.")
    bullet("Everything is documented. There's a full journal of every prompt we gave the AI, every bug it found, and every fix it applied. It's the real build story.")
    bullet("The AI hit real bugs on its own -- the same tricky platform behaviors that would trip up a human developer. The field-naming bug from the previous slide? The AI discovered that. The search bar crashing when placed in the wrong spot? The AI found and fixed that too.")
    pdf.ln(1)

    subheading("Right Column -- The Dashboard Story")
    body("\"People always want to know about the dashboard. Here's the story.\"")
    pdf.ln(1)
    bullet("First attempt: the AI reached for a popular charting library called Chart.js. It worked, but it was heavy -- a big dependency for what we needed.")
    bullet("Second attempt: I told the AI 'Replace the chart library with something simpler, just use basic web graphics.' One prompt. The AI rebuilt the entire analytics dashboard in 705 lines. The donut chart is literally just a circle with some math applied to it. No chart library at all.")
    bullet("Every chart, number, and bar in the dashboard is clickable. Click any visualization and it opens a table showing you the actual records behind that number. Simple filter callbacks -- nothing fancy.")
    bullet("The Kanban board is bigger at 1,117 lines, but it follows the same idea: standard, well-known tools doing standard things. Drag-and-drop library for card movement, a query library for data, a state library for UI, and Tailwind for styling.")
    pdf.ln(1)

    body("\"The real point isn't that an AI built it. The point is that Code Apps let you use "
         "standard web development tools, and standard tools are exactly what AI assistants are "
         "best at. The platform doesn't fight the tools, the tools don't fight the AI, and the AI "
         "can move fast because it's just writing a normal web app.\"")

    stage_direction("[Point to the stats row: 6 database tables, 18 build rounds, 705 dashboard "
                    "lines of code, 0 chart libraries, 1 deploy command. Let the numbers speak.]")

    divider()

    # ════════════════════════════════════════════════════════════════
    # SLIDE 5 TALK TRACK
    # ════════════════════════════════════════════════════════════════
    heading("Slide 5: Cool Things You Can Do")
    stage_direction("[Advance slide. Four pattern cards. Walk through each one briskly -- "
                    "you'll show these live in the demo.]")

    body("\"Here are four things from the app that I think are worth calling out. "
         "None of these are exotic or complicated. They're standard web development "
         "patterns. The point is that Code Apps let you use them -- the platform "
         "doesn't get in your way.\"")
    pdf.ln(1)

    subheading("Instant Search, Zero Extra Calls")
    body("\"The search bar -- Ctrl+K -- lets you find records across all six data tables "
         "instantly. It doesn't make any extra calls to the database. Here's why: as you "
         "navigate the app normally, the data you visit gets stored in memory. The search "
         "bar just reads that stored data. If you've visited the accounts page, those accounts "
         "are already in memory and searchable. No extra round trips.\"")
    pdf.ln(1)

    subheading("Quick-Create Buttons (27 Lines of Code)")
    body("\"The colored buttons along the top of every page are quick-create shortcuts. "
         "Tap 'work' and it opens a new action item form with the type already set to Work, "
         "on whatever page you're currently viewing. The entire mechanism behind this is 27 "
         "lines of code. One button sets a flag, the form component watches for that flag, "
         "and auto-opens. No complex wiring.\"")
    pdf.ln(1)

    subheading("Drag-and-Drop That Saves Automatically")
    body("\"The Kanban board has drag-and-drop between columns. Drag a card to the 'parking lot' "
         "column and the database immediately updates to flag that item as pinned. Drag it back "
         "out and the flag is removed. No save button, no confirmation dialog. The drag IS the save.\"")
    pdf.ln(1)

    subheading("Dark Mode That Works Inside the Platform")
    body("\"This one is subtle. The usual approach to dark mode -- detecting the user's system "
         "preference -- doesn't work in Code Apps because the Power Platform controls the outer "
         "page. You need to use a class-based approach instead, where your code explicitly adds "
         "a 'dark' marker to the page. Once you know this, it works perfectly. If you don't know "
         "this, you'll spend hours wondering why dark mode doesn't respond to your system settings.\"")

    stage_direction("[Keep this slide brisk. 60-90 seconds total. "
                    "The demo will show all four of these live.]")

    divider()

    # ════════════════════════════════════════════════════════════════
    # SLIDE 6 TALK TRACK
    # ════════════════════════════════════════════════════════════════
    heading("Slide 6: Let Me Show You")
    stage_direction("[Advance slide. Pause. Let the audience read the demo flow.]")

    body("\"Alright, enough slides. Let me show you the actual app.\"")
    pdf.ln(1)

    body("\"Here's what we'll walk through:\"")
    bullet("The analytics dashboard -- click any number to see the real records behind it.")
    bullet("The Kanban board -- I'll drag a card between columns so you can see the database update in real time.")
    bullet("Ctrl+K search -- find records across all data tables instantly, no extra database calls.")
    bullet("AI-powered action items -- I'll paste in some meeting notes and we'll watch the AI create structured action items automatically.")
    bullet("Dark mode -- one click, works inside the platform.")
    bullet("Under the hood -- I'll open dev tools so you can see the 77 scripts and the runtime we talked about.")
    bullet("Deploy -- two commands and it's live in Power Platform.")
    pdf.ln(2)

    stage_direction("[Switch to browser. App should already be loaded at the dashboard view. "
                    "If you need a fallback, have localhost:3001 running in a second tab.]")

    divider()

    # ════════════════════════════════════════════════════════════════
    # LIVE DEMO DETAILED BEATS
    # ════════════════════════════════════════════════════════════════
    heading("Live Demo: Detailed Beats")
    body("Below are the detailed beats for each demo segment. Times are approximate.")
    pdf.ln(2)

    subheading("Act 1: Dashboard (3 min)")
    stage_direction("[App is showing the analytics dashboard]")
    bullet("Point out the 4 summary cards at the top. Hover over one to show the tooltip with item names.")
    bullet("\"See this donut chart? No chart library. It's just a circle drawn with basic web graphics -- an SVG circle element with some math. 705 lines total for the entire dashboard.\"")
    bullet("Click a status row (e.g., 'In Progress') to open the drilldown. \"Every visualization on this page is clickable. Click a number and you see the actual records behind it.\"")
    bullet("Close the drilldown. Click a priority bar to show a different filter.")
    pdf.ln(1)

    subheading("Act 2: Kanban Board (4 min)")
    stage_direction("[Click 'My Board' in the sidebar]")
    bullet("Orient the audience: 4 columns -- parking lot (green, for pinned items), work (the main column), projects (purple), ideas (amber).")
    bullet("Hover over a card to show the floating toolbar: drag handle, color dots for priority, edit button, and pin icon.")
    bullet("Drag a card to the parking lot. \"That drag just saved to the database -- the item is now flagged as pinned in Dataverse.\" Drag it back out to unpin.")
    bullet("Click the filter buttons in the work column header (A/W/P/L for All/Work/Personal/Learning). \"Watch the column title, icon, and color change with each filter.\"")
    bullet("Click the edit pencil on a card. \"It opens the correct form based on what type of item it is -- action item, project, idea, or meeting summary.\"")
    pdf.ln(1)

    subheading("Act 3: Search (2 min)")
    stage_direction("[Press Ctrl+K]")
    bullet("Type a search term that spans multiple data types (e.g., a company name).")
    bullet("\"Notice how fast that is. Zero database calls. It's reading data that was already loaded when you visited those pages. The search just looks through what's already in memory.\"")
    bullet("Arrow down and press Enter to navigate to a result. Show that it takes you to the right page.")
    pdf.ln(1)

    subheading("Act 4: AI-Powered Action Items (5 min)")
    stage_direction("[Navigate to Meeting Summaries. Open a meeting summary that has notes, "
                    "or create one and paste in sample meeting notes.]")
    bullet("Click 'Extract Action Items with AI.'")
    bullet("\"This sends the meeting text to Azure OpenAI. The AI reads the notes and returns structured data -- a name, priority level, due date, and description for each action item.\"")
    bullet("Show the results. Point out that the AI already assigned the right priority codes that match the database -- it's not just text, it's ready to save.")
    bullet("Save the action items. Navigate to the action items list to show they're there.")
    bullet("\"The entire AI integration is about 80 lines of code. If the AI service isn't configured, the button simply doesn't appear -- no error, no broken feature.\"")
    pdf.ln(1)

    subheading("Act 5: Dark Mode + Under the Hood (2 min)")
    stage_direction("[Click the moon/sun icon in the sidebar footer]")
    bullet("Toggle dark mode. \"One click. It remembers your choice. And it works inside the Power Platform iframe because we used a class-based approach instead of media queries.\"")
    bullet("Open browser dev tools (F12). Show the Elements panel.")
    bullet("\"Remember slide 2? Here's what we were talking about. You can see the nested iframes, the scripts, the feature switches. Your React app is in here -- but it's a small piece inside all of this.\"")
    pdf.ln(1)

    subheading("Act 6: Deploy (1 min)")
    stage_direction("[Switch to terminal -- or just narrate this if time is short]")
    code("npm run build && pac code push")
    bullet("\"That's the whole deploy. The first command builds your app. The second command pushes it to Power Platform. Same security, same governance, same compliance as any Canvas App. No separate web hosting needed.\"")
    pdf.ln(1)

    subheading("Act 7: Wrap-Up (1 min)")
    body("\"So here's what we showed today: a modern web app, built entirely by an AI assistant, "
         "running inside the Power Platform with full database integration, AI-powered features, "
         "drag-and-drop, instant search, dark mode, and a one-command deploy. "
         "Code Apps give your developers a way into Power Platform without giving up "
         "any of the tools they already know and love.\"")
    pdf.ln(2)

    divider()

    # ════════════════════════════════════════════════════════════════
    # RECOVERY PLAYS
    # ════════════════════════════════════════════════════════════════
    heading("Recovery Plays")
    body("Things that might go wrong during the live demo and how to handle them.")
    pdf.ln(1)

    subheading("Database is Slow")
    body("If creates or updates take a moment: \"The database is in the cloud -- sometimes there's "
         "a brief delay. You'd see a loading spinner in production. The app automatically picks up "
         "the change as soon as it lands.\"")

    subheading("AI Extraction Fails")
    body("If the AI service times out or gives unexpected results: \"This is a live API call to "
         "Azure OpenAI -- sometimes there's a hiccup. The app shows a clear error message when that "
         "happens. In production you'd add retry logic. The point here is the integration pattern, "
         "not the uptime of my demo API key.\"")

    subheading("\"Is this a Canvas App or a Model-Driven App?\"")
    body("\"Neither, and kind of both. Code Apps are a third option. You write standard web code "
         "(React, TypeScript), you deploy to Power Platform, and you get the same database, the same "
         "governance, and the same security model. The platform even classifies it as a Canvas App "
         "internally. But your source code is standard web development.\"")

    subheading("\"Can I use this with Copilot Studio?\"")
    body("\"Yes. We showed that the platform has a Copilot panel pre-wired into every Code App -- "
         "we saw 185 references to it in the page code. And the AI extraction feature in this demo "
         "talks to Azure OpenAI directly. Both approaches work.\"")

    subheading("\"How do I get started?\"")
    body("\"My colleague covered that in the intro talk. Three commands: pac code init to create "
         "the project, pac code add-data-source to connect your database tables, and npm run dev "
         "to start coding. The getting-started story is straightforward. What I showed you today "
         "is what happens when you push past hello world.\"")

    pdf_path = os.path.join(OUT_DIR, "code-apps-under-the-hood-talk-track.pdf")
    pdf.output(pdf_path)
    print(f"PDF saved: {pdf_path}")
    return pdf_path


# ── Main ────────────────────────────────────────────────────────────
if __name__ == "__main__":
    build_pptx()
    build_pdf()
    print("Done!")
